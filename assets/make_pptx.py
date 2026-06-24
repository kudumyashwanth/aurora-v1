from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
A=os.path.abspath("assets")
BG=RGBColor(0x0d,0x1b,0x2a); FG=RGBColor(0xe6,0xed,0xf3); AC=RGBColor(0x00,0xd3,0xa7)
SUB=RGBColor(0x9f,0xb3,0xc8); AM=RGBColor(0xff,0xb3,0x47); GR=RGBColor(0x7c,0xfc,0x98)
prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
SW,SH=prs.slide_width,prs.slide_height
def slide():
    s=prs.slides.add_slide(prs.slide_layouts[6])
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background()
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def tb(s,x,y,w,h,txt,sz,col,bold=True,al=PP_ALIGN.LEFT,font="Calibri"):
    b=s.shapes.add_textbox(x,y,w,h); tf=b.text_frame; tf.word_wrap=True
    lines=txt.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=al
        run=p.add_run(); run.text=ln; f=run.font; f.size=Pt(sz); f.bold=bold; f.color.rgb=col; f.name=font
    return b
def bar(s):
    r=s.shapes.add_shape(1,0,0,I(0.18),SH); r.fill.solid(); r.fill.fore_color.rgb=AC; r.line.fill.background()
def pic(s,path,x,y,w):
    return s.shapes.add_picture(path,x,y,width=w)

# 1 TITLE
s=slide(); bar(s)
tb(s,I(0.9),I(2.0),I(11.5),I(0.5),"RISC-V  ·  TENSOR ACCELERATOR  ·  AXI4  ·  SKY130B  ·  RTL → GDSII",16,AC)
tb(s,I(0.85),I(2.5),I(11.8),I(1.3),"Aurora v1 — AI SoC",54,FG)
tb(s,I(0.9),I(3.9),I(11.5),I(1.2),"A complete System-on-Chip taken from RTL to a signed-off GDSII layout —\nopen-source tools only, on a single 23 GB laptop.",22,SUB,False)
tb(s,I(0.9),I(6.4),I(11.5),I(0.6),"Yashwanth Kudum   ·   github.com/kudumyashwanth/aurora-v1",16,SUB,False)

# 2 OVERVIEW
s=slide(); bar(s); tb(s,I(0.6),I(0.4),I(12),I(0.9),"What is Aurora v1?",34,AC)
tb(s,I(0.7),I(1.6),I(7.0),I(4.5),
 "Multi-clock AI SoC: a 64-bit RISC-V core + a hardware matrix-multiply\naccelerator over a coherent-terminated AXI4 fabric.\n\n"
 "•  Rocket RV64IMAC CPU tile — hardened macro, 33 MHz\n"
 "•  4×4 systolic tensor engine (int16) — hardened macro, 50 MHz\n"
 "•  8×8 AXI4 crossbar + dual async clock-domain bridges\n"
 "•  UART · GPIO · timer · interrupt ctrl · boot ROM · SRAM\n\n"
 "Carried through the ENTIRE flow: RTL → lint → sim → formal →\nsynthesis → floorplan → place → CTS → route → DRC → LVS → GDSII.",18,FG,False)
for i,(k,v) in enumerate([("RTL→GDSII","complete"),("Devices","137,562"),("Die","41.25 mm²"),("Clocks","33 / 50 MHz")]):
    y=I(1.7+i*1.25)
    tb(s,I(8.2),y,I(4.5),I(0.5),v,30,AC); tb(s,I(8.2),y+I(0.62),I(4.5),I(0.4),k,15,SUB,False)

# 3 ARCHITECTURE
s=slide(); bar(s); tb(s,I(0.6),I(0.35),I(12),I(0.8),"Architecture",32,AC)
pic(s,A+"/architecture.png",I(1.0),I(1.25),I(11.3))

# 4 LAYOUT
s=slide(); bar(s); tb(s,I(0.6),I(0.35),I(12),I(0.8),"Full-Chip Layout",32,AC)
p=pic(s,A+"/chip_floorplan.png",I(2.7),I(1.2),I(8.0))
tb(s,I(0.45),I(2.3),I(2.3),I(3),"Rendered from\nthe actual\nplacement\ndatabase —\n130.5k glue\ncells + 2\nhard macros.",14,SUB,False)

# 5 RESULTS
s=slide(); bar(s); tb(s,I(0.6),I(0.35),I(12),I(0.8),"Sign-off Results",32,AC)
pic(s,A+"/results_dashboard.png",I(1.1),I(1.3),I(11.1))

# 6 BOOT
s=slide(); bar(s); tb(s,I(0.6),I(0.35),I(12),I(0.8),"It Boots — Functional Verification",32,AC)
pic(s,A+"/boot_terminal.png",I(2.6),I(1.35),I(8.1))
tb(s,I(0.45),I(2.4),I(2.2),I(3),"Full-SoC sim:\nRocket boots\nthrough the\nreal fabric,\ndrives the\ntensor core,\nreads result\nover AXI.\nZero traps.",14,SUB,False)

# 7 HIGHLIGHTS
s=slide(); bar(s); tb(s,I(0.6),I(0.4),I(12),I(0.9),"Engineering Highlights",34,AC)
tb(s,I(0.7),I(1.7),I(11.8),I(5),
 "■  Cracked a routing-congestion wall on the dense systolic datapath —\n     diagnosed a metal-1 short-storm, beat it with a soft layer-derate.\n\n"
 "■  Hardened a RISC-V tile as an AXI macro — re-wrapped Rocket's\n     TileLink-C port through a cache cork + TL-to-AXI4 bridge (Chisel).\n\n"
 "■  Closed multi-clock timing — moved the FPU/divider-bound CPU path\n     to RV64IMAC @ 33 MHz; accelerator + fabric @ 50 MHz, async-bridged.\n\n"
 "■  Full-chip PDN + LVS integration — root-caused power-grid fragmentation\n     to a dual-layer-power macro; fixed with a PDN keep-out → clean device-match.\n\n"
 "■  Memory-bounded P&R — entire flow capped at 18 GB to survive a 23 GB host.",18,FG,False)

# 8 CLOSING
s=slide(); bar(s)
tb(s,I(0.85),I(2.3),I(11.8),I(1.2),"RTL to silicon — done.",46,FG)
tb(s,I(0.9),I(3.7),I(11.5),I(1.4),"A full RISC-V + tensor AI SoC, signed off to GDSII on commodity\nhardware with open tools. Timing MET · DRC 1 · LVS device-match.",22,SUB,False)
tb(s,I(0.9),I(5.9),I(11.5),I(0.6),"github.com/kudumyashwanth/aurora-v1   ·   IP available for licensing",17,AC)
prs.save("deliverables/Aurora_v1.pptx")
print("PPTX slides:",len(prs.slides._sldIdLst))
